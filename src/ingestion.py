from __future__ import annotations

import hashlib
import logging
from pathlib import Path

from docx import Document as DocxDocument
from langchain_core.documents import Document
from odf import opendocument, text as odf_text
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
import pytesseract
from pytesseract import TesseractNotFoundError

from .config import (
    SUPPORTED_EXTENSIONS,
    SUPPORTED_IMAGE_EXTENSIONS,
    SUPPORTED_OFFICE_EXTENSIONS,
    SUPPORTED_PDF_EXTENSIONS,
    SUPPORTED_TEXT_EXTENSIONS,
)

logger = logging.getLogger(__name__)


def collect_supported_files(folder: str | Path) -> list[Path]:
    root = Path(folder)
    if not root.exists():
        return []

    files = [
        path
        for path in root.rglob("*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    return sorted(files)


def calculate_file_hash(path: str | Path) -> str:
    """Calculate SHA256 hash of file contents for duplicate detection."""
    file_path = Path(path)
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()


def load_documents(paths: list[str | Path]) -> list[Document]:
    documents: list[Document] = []
    for path in paths:
        try:
            documents.extend(load_document(path))
        except Exception as error:
            logger.error(f"Failed to load {path}: {error}")
    return documents


def load_document(path: str | Path) -> list[Document]:
    file_path = Path(path)
    suffix = file_path.suffix.lower()

    if suffix in SUPPORTED_PDF_EXTENSIONS:
        return _load_pdf(file_path)
    if suffix in SUPPORTED_TEXT_EXTENSIONS:
        return _load_text(file_path)
    if suffix in SUPPORTED_IMAGE_EXTENSIONS:
        return _load_image(file_path)
    if suffix in SUPPORTED_OFFICE_EXTENSIONS:
        return _load_office(file_path)

    raise ValueError(f"Unsupported file type: {file_path.suffix}")


def _load_pdf(path: Path) -> list[Document]:
    reader = PdfReader(str(path))
    documents: list[Document] = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue

        documents.append(
            Document(
                page_content=text,
                metadata={
                    "source": path.name,
                    "source_path": str(path),
                    "type": "pdf",
                    "page": page_number,
                },
            )
        )

    return documents


def _load_text(path: Path) -> list[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        return []

    return [
        Document(
            page_content=text,
            metadata={
                "source": path.name,
                "source_path": str(path),
                "type": "text",
                "page": 1,
            },
        )
    ]


def _load_image(path: Path) -> list[Document]:
    try:
        with Image.open(path) as image:
            text = pytesseract.image_to_string(image).strip()
    except (TesseractNotFoundError, Exception) as error:
        logger.warning(
            f"OCR failed for {path.name}: {error}. Storing image metadata only."
        )
        text = f"[OCR failed] Image file: {path.name}"

    if not text:
        return []

    return [
        Document(
            page_content=text,
            metadata={
                "source": path.name,
                "source_path": str(path),
                "type": "image",
                "page": 1,
            },
        )
    ]


def _load_office(path: Path) -> list[Document]:
    suffix = path.suffix.lower()

    if suffix == ".docx":
        return _load_docx(path)
    elif suffix == ".pptx":
        return _load_pptx(path)
    elif suffix == ".xlsx":
        return _load_xlsx(path)
    elif suffix == ".odt":
        return _load_odt(path)
    else:
        raise ValueError(f"Unsupported office format: {suffix}")


def _load_docx(path: Path) -> list[Document]:
    doc = DocxDocument(str(path))
    paragraph_lines = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]

    # Speaker notes and structured content are often stored in DOCX tables.
    table_lines: list[str] = []
    for table in doc.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                table_lines.append(" | ".join(row_cells))

    text_parts = paragraph_lines + table_lines
    text = "\n".join(text_parts).strip()

    if not text:
        return []

    return [
        Document(
            page_content=text,
            metadata={
                "source": path.name,
                "source_path": str(path),
                "type": "docx",
                "page": 1,
            },
        )
    ]


def _load_pptx(path: Path) -> list[Document]:
    presentation = Presentation(str(path))
    documents: list[Document] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        text = "\n".join(text_parts).strip()
        if not text:
            continue

        documents.append(
            Document(
                page_content=text,
                metadata={
                    "source": path.name,
                    "source_path": str(path),
                    "type": "pptx",
                    "page": slide_number,
                },
            )
        )

    return documents


def _load_xlsx(path: Path) -> list[Document]:
    workbook = load_workbook(str(path))
    documents: list[Document] = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows: list[list[str]] = []

        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)

        text = "\n".join(rows).strip()
        if not text:
            continue

        documents.append(
            Document(
                page_content=text,
                metadata={
                    "source": path.name,
                    "source_path": str(path),
                    "type": "xlsx",
                    "sheet": sheet_name,
                    "page": 1,
                },
            )
        )

    return documents


def _load_odt(path: Path) -> list[Document]:
    doc = opendocument.load(str(path))
    text_parts: list[str] = []

    for paragraph in doc.getElementsByType(odf_text.P):
        text = "".join(
            node.data for node in paragraph.childNodes if hasattr(node, "data")
        ).strip()
        if text:
            text_parts.append(text)

    text = "\n".join(text_parts).strip()

    if not text:
        return []

    return [
        Document(
            page_content=text,
            metadata={
                "source": path.name,
                "source_path": str(path),
                "type": "odt",
                "page": 1,
            },
        )
    ]
